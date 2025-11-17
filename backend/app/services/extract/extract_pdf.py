from __future__ import annotations

import csv
import io
import logging
import re
from typing import Optional

from pypdf import PdfReader 
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd


logger = logging.getLogger(__name__)

def extract_text_from_pdf(data: bytes) -> str:
    """Extrai texto de arquivos PDF."""
    try:
        buf = io.BytesIO(data)
        reader = PdfReader(buf)
        texts: list[str] = []
        for page in reader.pages:
            try:
                txt = page.extract_text() or ""
                if txt.strip():
                    texts.append(txt.strip())
            except Exception:
                continue
        return "\n".join(texts).strip()
    except Exception:
        logger.exception("Falha ao extrair texto de PDF")
        return ""


def extract_text_from_docx(data: bytes) -> str:
    """Extrai texto de arquivos Word (.docx)."""
    try:
        buf = io.BytesIO(data)
        document = DocxDocument(buf)
        parts: list[str] = []
        
        for p in document.paragraphs:
            text = (p.text or "").strip()
            if text:
                parts.append(text)
        
        for table in getattr(document, 'tables', []) or []:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if (cell.text or '').strip()]
                if row_text:
                    parts.append(" | ".join(row_text))
        
        return "\n".join(parts).strip()
    except Exception:
        logger.exception("Falha ao extrair texto de DOCX")
        return ""


def extract_text_from_excel(data: bytes) -> str:
    """Extrai texto de arquivos Excel (.xlsx, .xls)."""
    try:
        buf = io.BytesIO(data)
        try:
            df = pd.read_excel(buf, sheet_name=None, engine='openpyxl')
        except Exception:
            buf.seek(0)
            df = pd.read_excel(buf, sheet_name=None, engine='xlrd')
        
        parts: list[str] = []
        for sheet_name, sheet_df in df.items():
            parts.append(f"=== {sheet_name} ===")
            sheet_text = sheet_df.fillna('').astype(str)
            for row in sheet_text.iterrows():
                row_text = [str(cell).strip() for cell in row[1] if str(cell).strip()]
                if row_text:
                    parts.append(" | ".join(row_text))
        return "\n".join(parts).strip()
    except Exception:
        logger.exception("Falha ao extrair texto de Excel")
        return ""


def extract_text_from_csv(data: bytes) -> str:
    """Extrai texto de arquivos CSV."""
    try:
        text = data.decode('utf-8')
        parts: list[str] = []
        reader = csv.reader(io.StringIO(text))
        for row in reader:
            row_text = [cell.strip() for cell in row if cell.strip()]
            if row_text:
                parts.append(" | ".join(row_text))
        return "\n".join(parts).strip()
    except UnicodeDecodeError:
        try:
            text = data.decode('latin-1')
            parts: list[str] = []
            reader = csv.reader(io.StringIO(text))
            for row in reader:
                row_text = [cell.strip() for cell in row if cell.strip()]
                if row_text:
                    parts.append(" | ".join(row_text))
            return "\n".join(parts).strip()
        except Exception:
            logger.exception("Falha ao extrair texto de CSV")
            return ""
    except Exception:
        logger.exception("Falha ao extrair texto de CSV")
        return ""


def extract_text_from_pptx(data: bytes) -> str:
    """Extrai texto de arquivos PowerPoint (.pptx)."""
    try:
        buf = io.BytesIO(data)
        presentation = Presentation(buf)
        parts: list[str] = []
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        parts.append(text)
                        
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            parts.append(" | ".join(row_text))
        
        return "\n".join(parts).strip()
    except Exception:
        logger.exception("Falha ao extrair texto de PPTX")
        return ""


def extract_text_from_doc(data: bytes) -> str:
    """Extrai texto básico de arquivos Word legado (.doc)."""
    try:
        text = data.decode('latin-1', errors='ignore')
        
        text = re.sub(r'Microsoft Word.*?Document.*?Word\.Document\.\d+.*?Caolan\d*', '', text)
        text = re.sub(r'MSWordDoc.*?Word\.Document\.\d+', '', text)
        text = re.sub(r'Root Entry.*?WordDocument', '', text)
        text = re.sub(r'CompObj.*?SummaryInformation', '', text)
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', ' ', text)
        
        readable_parts = re.findall(r'[A-Za-z0-9À-ÿ\s.,;:!?()"\'\-–—]{8,}', text)
        
        if readable_parts:
            extracted = ' '.join(readable_parts)
            extracted = re.sub(r'[^\w\s.,;:!?()\'"À-ÿ\-–—]{3,}', ' ', extracted)
            extracted = re.sub(r'\s+', ' ', extracted).strip()
            
            words = re.findall(r'\b[A-Za-zÀ-ÿ]{3,}\b', extracted)
            if len(words) >= 5:
                return extracted
        
        return ""
    except Exception:
        logger.exception("Falha ao extrair texto de DOC")
        return ""


def extract_text_from_txt(data: bytes) -> str:
    """Extrai texto de arquivos de texto simples."""
    try:
        return data.decode('utf-8').strip()
    except UnicodeDecodeError:
        try:
            return data.decode('latin-1').strip()
        except Exception:
            logger.exception("Falha ao extrair texto de arquivo TXT")
            return ""
    except Exception:
        logger.exception("Falha ao extrair texto de arquivo TXT")
        return ""


def extract_text_from_document(data: bytes, mime_type: Optional[str] = None, filename: Optional[str] = None) -> str:
    """
    Extrai texto de diversos tipos de documentos.
    
    Suporta:
    - PDF (.pdf)
    - Word (.docx, .doc)
    - PowerPoint (.pptx)
    - Excel (.xlsx, .xls)
    - CSV (.csv)
    - Texto (.txt, .text, .md)
    """
    mt = (mime_type or '').lower()
    name = (filename or '').lower()
    
    if mt.endswith('pdf') or name.endswith('.pdf'):
        return extract_text_from_pdf(data)
    
    if ('presentation' in mt or 'powerpoint' in mt or name.endswith('.pptx')):
        return extract_text_from_pptx(data)
    
    if (('word' in mt and 'document' in mt) or 'wordprocessingml' in mt or name.endswith('.docx')):
        return extract_text_from_docx(data)
    
    if (('msword' in mt and 'document' not in mt) or name.endswith('.doc')):
        return extract_text_from_doc(data)
    
    if ('excel' in mt or 'spreadsheet' in mt or name.endswith('.xlsx') or name.endswith('.xls')):
        return extract_text_from_excel(data)
    
    if 'csv' in mt or name.endswith('.csv'):
        return extract_text_from_csv(data)
    
    if ('text' in mt or name.endswith('.txt') or name.endswith('.text') or name.endswith('.md')):
        return extract_text_from_txt(data)
    
    logger.info("Tipo de documento não suportado para extração: mime=%s, name=%s", mime_type, filename)
    return ""
